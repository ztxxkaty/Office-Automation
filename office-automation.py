## -------- 0、基础设置 -------- ##
# 下载库
# pip install pandas # 在命令行窗口中输入下载

# 载入库
import pandas as pd
import matplotlib.pyplot as plt
import docx
from pptx import Presentation
from pptx.util import Inches

## -------- 1、筛选符合条件的候选人 -------- ##
# excel
excel_name = 'candidates.xlsx'
figure_name = "xinzi.png"

# 读取excel文件
candidates_df = pd.read_excel(excel_name)

# 筛选符合条件的候选人
filtered_candidates = candidates_df[candidates_df['工作经验（年）']>3] # 筛选条件：工作经验>3年

# 数据分析
df = pd.DataFrame(filtered_candidates)

avg_age = df['年龄'].mean()

plt.bar(df['姓名'],df['期望薪资'], color='skyblue')
plt.rcParams['font.sans-serif'] = ['SimHei'] # 设置字符
plt.title('期望薪资')
# plt.show()
plt.savefig(figure_name)

## -------- 2、生成面试通知 -------- ##
# word
photo_name = 'photo.png'
interview_details = {
    '时间': '2023年10月15日上午10点',
    '地点': '公司总部大楼3楼会议室'
}

for index, candidate in filtered_candidates.iterrows():
    doc = docx.Document()
    
    doc.styles['Normal'].font.size=docx.shared.Pt(18)
    doc.styles['Normal'].font.name='Times New Roman' # 设置英文字符的字体
    doc.styles['Normal']._element.rPr.rFonts.set(docx.oxml.ns.qn('w:eastAsia'),'微软雅黑') # 设置中文字符的字体
    # doc.styles['Normal]._element.rPr.rFonts.st(qn('w:eastAsia
    doc.add_heading('面试邀请',0)
    doc.add_paragraph(f"尊敬的{candidate['姓名']}：")
    doc.add_paragraph("谢谢您选择向我们投递简历。")
    doc.add_paragraph("我们诚挚地邀请您参加面试，以下是面试的详细信息：")

    for key, value in interview_details.items():
        doc.add_paragraph(f"{key}：{value}")
    
    # 插入图表
    doc.add_picture(photo_name, width=Inches(5))
    
    # 保存文档
    doc_path = f"{candidate['姓名']}_面试通知.docx"
    doc.save(doc_path)
    
## -------- 3、生成汇报ppt -------- ##
# ppt
ppt_name = '候选人情况汇总.pptx'
ppt_title = "候选人情况汇总"
image_Folder = 'image'

ppt = Presentation()

# 添加标题页幻灯片
title_slide_layout = ppt.slide_layouts[0] # 0--标题页
title_slide = ppt.slides.add_slide(title_slide_layout)
title = title_slide.shapes.title
subtitle = title_slide.placeholders[1]
title.text = ppt_title
subtitle.text = f"符合招聘要求的候选人共有{len(filtered_candidates)}人，平均年龄{avg_age}岁"  # 注意：这里应该插入实际的日期

# 添加薪资分布
slide_layout = ppt.slide_layouts[6]
slide = ppt.slides.add_slide(slide_layout)
left = Inches(1.5)
top = Inches(1.5)
pic = slide.shapes.add_picture(figure_name, left, top)

# 为每位候选人添加一个幻灯片页
for index, row in filtered_candidates.iterrows():
    slide_layout = ppt.slide_layouts[1]
    slide = ppt.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = row['姓名']
    
    content_placeholder = slide.placeholders[1]
    text_frame = content_placeholder.text_frame
    text_frame.text = (
        f"工作经验：{row['工作经验（年）']}年\n"
        f"专业：{row['专业']}\n"
        f"年龄：{row['年龄']}岁\n" 
        f"申请岗位：{row['申请岗位']}\n" 
        )
    
    image_name = f"{image_Folder}\{row['姓名']}.jpg"
    left = Inches(7)
    top = Inches(1)
    # width = Inches(1)
    height = Inches(2)
    # pic = slide.shapes.add_picture(image_name, left, top, width, height)
    pic = slide.shapes.add_picture(image_name, left, top, height=height)

# 保存PPT
ppt.save(ppt_name)
